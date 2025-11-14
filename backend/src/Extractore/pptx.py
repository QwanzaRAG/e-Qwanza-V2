import warnings
warnings.simplefilter(action='ignore', category=FutureWarning)

import io
import os
import shutil
import pandas as pd
import numpy as np

from enum import Enum
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.util import Inches, Cm, Pt
from pptx.oxml.xmlchemy import OxmlElement
from pptx.dml.color import ColorFormat, RGBColor
from typing import List
from .common_functions import *

# from ..utils.common_functions import *
from .image import extract_text_from_ocr, ImageChartDataExtractor
class Entity:
    """
    Represents an entity with chart type, text, position, and size.
    """
    def __init__(self, chart_type: str, text: str, left: int, top: int, width: int, height: int) -> None:
        """
        Initializes the Entity object.

        Args:
            chart_type (str): The type of chart.
            text (str): The text associated with the entity.
            left (int): The left position of the entity.
            top (int): The top position of the entity.
            width (int): The width of the entity.
            height (int): The height of the entity.
        """
        self.chart_type = chart_type
        self.text = text
        self.left = left
        self.top = top
        self.width = width
        self.height = height
        
    def __repr__(self): return repr(self.__dict__)
          
class Slide:
    """
    Represents a slide with slide number, title, text, entities.
    """
    def __init__(self, slide_number: int, slide_title: str, slide_text: str, entities: List[Entity]) -> None:
        """
        Initializes the Slide object.

        Args:
            slide_number (int): The number of the slide.
            slide_title (str): The title of the slide.
            slide_text (str): The text content of the slide.
            entities (List[Entity]): The list of entities on the slide.
        """
        self.slide_number = slide_number
        self.slide_title = slide_title
        self.slide_text = slide_text
        self.entities = entities
        self.similarity = 0
    
    def __repr__(self): return repr(self.__dict__)
        
        
class PPTExtractor():
    def __init__(self, file_path, extraction_method: str = "slide", ocr_engine: str = "tesseract",extract_from_image: bool = True,chart_from_image : bool = True) -> None:
        """
        Initializes a PPTExtractor object.

        Args:
        - file_path (str): The path to the PPT file.
        - extraction_method (str, optional): The method to extract content from slides. Defaults to "slide".
        - ocr_engine (str, optional): The OCR engine to use for image text extraction. Defaults to "tesseract".
        """
        self.file_path = file_path
        self.extraction_method = extraction_method
        self.ocr_engine = ocr_engine
        self.slides = []
        self.extract_from_image = extract_from_image
        self.chart_from_image = chart_from_image

    def extract(self,maintain_order : bool = False):
        """
        Extracts content from the PPT file.

        Raises:
        - Exception: If the PPT file is not found.
        """
        if not os.path.isfile(self.file_path):
            raise Exception("PPT File not Found")

        presentation = Presentation(self.file_path)
        self.title = presentation.core_properties.title
        self.author = presentation.core_properties.author
        self.subject = presentation.core_properties.subject
        self.keywords = presentation.core_properties.keywords
        self.last_modified_by = presentation.core_properties.last_modified_by
        self.created = presentation.core_properties.created
        self.modified = presentation.core_properties.modified
        self.slide_height = presentation.slide_height
        self.slide_width = presentation.slide_width
        self.entities = []
        self.deplot = ImageChartDataExtractor()
        
        for slide_number, slide in enumerate(presentation.slides):
            try:
                slide_number += 1
                entities = []
                slide_wise_text = ''
                slide_wise_table = ''
                slide_wise_chart = ''
                slide_wise_ocr = ''

                slide_title = ""
                try:
                    slide_title += slide.shapes.title.text
                except:
                    pass
                
                
                slide_text = "Slide Number : "+str(slide_number)
                slide_text += "\nSlide Title : "+str(slide_title)
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                slide_wise_text += run.text
                                slide_text += "\nSlide Text : "+run.text
                                entities.append(Entity("text", run.text, shape.left, shape.top, shape.width, shape.height))
                    
                    # Extract table from shapes
                    if shape.has_table:
                        table = shape.table
                        table_str = convert_pptx_table_to_prettytable(table)
                        slide_wise_table += " \n " + table_str + " \n "
                        slide_text += "\nSlide Table : "+table_str
                        entities.append(Entity("table", table_str, shape.left, shape.top, shape.width, shape.height))
                    
                    # Extract chart from shapes
                    if shape.has_chart:
                        chart = shape.chart
                        chart_data_df = pd.read_excel(chart.part.chart_workbook.xlsx_part.blob).dropna(axis=1, how='all').dropna(axis=0, how='all')
                        chart_str = format_dataframe_to_prettytables(chart_data_df)
                        chart_text = ''
                        if chart.has_title:
                            chart_text += "Chart Title : " + chart.chart_title.text_frame.text
                        chart_text += "\nChart Type : " + str(chart.chart_type) + "\n"
                        chart_text += chart_str
                        slide_wise_chart += " \n " + chart_text + " \n "
                        slide_text += "\nSlide Chart : "+chart_text
                        entities.append(Entity("chart", chart_text, shape.left, shape.top, shape.width, shape.height))
                    
                    # Extract OCR text from images
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and self.extract_from_image:
                        if self.chart_from_image:
                            text = self.deplot.extract_chart_data_from_image(Image.open(io.BytesIO(shape.image.blob))) 
                            # text = self.deplot.extract_chart_data_from_image(io.BytesIO(shape.image.blob)) 
                        else:
                            text = extract_text_from_ocr(Image.open(io.BytesIO(shape.image.blob)))
                            # text = extract_text_from_ocr(io.BytesIO(shape.image.blob))
                        slide_wise_ocr += text + " \n "
                        slide_text += "\nSlide OCR : "+text
                        entities.append(Entity("image", text, shape.left, shape.top, shape.width, shape.height))

                if not maintain_order:
                    slide_text = f"""
                    Slide Number {slide_number}
                    Slide Title : {slide_title}
                    Slide Text : {slide_wise_text}
                    Slide Table : 
                    {slide_wise_table}
                    Slide Charts Data : 
                    {slide_wise_chart}
                    Slide Image OCR Text : 
                    {slide_wise_ocr}
                    """
                print("slide_Number: ", slide_number)
                print("slide_Number: ", slide_text)
                self.slides.append(Slide(slide_number, slide_title, slide_text, entities))
            except Exception as e:
                print("Ignoring slide ", slide_number, " Error ", e)
        
        